import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

# --- Streamlit Page Setup ---
st.set_page_config(page_title="Control Chart Maker (QA OMAC Tools)", layout="wide")
st.markdown("<h1 style='text-align:left;'>Control chart maker (QA OMAC Tools)</h1>", unsafe_allow_html=True)
st.markdown("---")

# --- Sidebar ---
st.sidebar.header("Upload & Setup")
uploaded_file = st.sidebar.file_uploader("Upload an Excel workbook (.xlsx)", type=["xlsx"]) 

if uploaded_file is not None:
    try:
        xl = pd.ExcelFile(uploaded_file)
        sheets = xl.sheet_names
    except Exception as e:
        st.sidebar.error(f"Unable to read workbook: {e}")
        st.stop()

    sheet_name = st.sidebar.selectbox("Select sheet that contains data", options=sheets)
    if sheet_name:
        df = pd.read_excel(xl, sheet_name=sheet_name, header=0)
        st.write("**Preview of selected sheet**")
        st.dataframe(df.head())

        # --- First column is x-axis ---
        time_col = df.columns[0]
        st.sidebar.write(f"Time / Batch column detected: **{time_col}**")

        # --- Select parameter columns ---
        st.sidebar.subheader("Parameters for Control Chart")
        param_cols = st.sidebar.multiselect(
            "Select one or more parameter columns",
            options=[c for c in df.columns if c != time_col],
            default=[c for c in df.columns if c != time_col]
        )

        if len(param_cols) == 0:
            st.info("Select at least one parameter column from the sidebar to proceed.")
        else:
            if st.sidebar.button("Process & Generate Outputs"):
                working = df.copy()

                # --- Safe numeric conversion ---
                for col in param_cols:
                    working[col] = (
                        working[col]
                        .astype(str)
                        .str.replace(',', '', regex=False)
                        .str.strip()
                        .replace(['-', '–', '—', 'N/A', 'na', '', None], np.nan)
                    )
                    working[col] = pd.to_numeric(working[col], errors='coerce')

                # --- Prepare Excel output ---
                out_xlsx = BytesIO()
                writer = pd.ExcelWriter(out_xlsx, engine='openpyxl')

                # --- Prepare PPTX ---
                prs = Presentation()
                slide_layout = prs.slide_layouts[5]

                # --- Process each parameter ---
                for col in param_cols:
                    col_values = working[col].dropna()
                    if col_values.empty:
                        st.warning(f"No numeric data found for '{col}'. Skipping.")
                        continue

                    # Calculate statistics
                    CL = col_values.mean()
                    MR = col_values.diff().abs()
                    MRbar = MR[1:].mean() if len(MR) > 1 else 0
                    d2 = 1.128
                    UCL = CL + 3 * (MRbar / d2)
                    LCL = max(0, CL - 3 * (MRbar / d2))

                    # Append new columns
                    working[f'{col}_CL'] = CL
                    working[f'{col}_MR'] = MR
                    working[f'{col}_MRbar'] = MRbar
                    working[f'{col}_UCL'] = UCL
                    working[f'{col}_LCL'] = LCL

                    # --- Create chart ---
                    fig, ax = plt.subplots(figsize=(10,4))
                    ax.plot(working[time_col], col_values, marker='o', label=col)
                    ax.axhline(CL, linestyle='--', label='CL')
                    ax.axhline(UCL, color='r', linestyle='--', label='UCL')
                    ax.axhline(LCL, color='r', linestyle='--', label='LCL')
                    ax.set_title(f'Control Chart - {col}')
                    ax.set_xlabel(time_col)
                    ax.set_ylabel(col)
                    ax.legend()
                    plt.xticks(rotation=45)
                    plt.tight_layout()

                    # Save chart as image
                    img_bytes = BytesIO()
                    fig.savefig(img_bytes, format='png', dpi=150)
                    img_bytes.seek(0)
                    plt.close(fig)

                    # Add slide
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(img_bytes, Inches(0.5), Inches(0.7), width=Inches(9))

                # --- Save Excel ---
                working.to_excel(writer, sheet_name=sheet_name, index=False)
                writer.close()
                out_xlsx.seek(0)

                # --- Save PPTX ---
                pptx_io = BytesIO()
                prs.save(pptx_io)
                pptx_io.seek(0)

                # --- Streamlit download buttons ---
                st.success('Processing complete!')
                st.download_button('Download updated Excel', data=out_xlsx,
                                   file_name='control_chart_output.xlsx',
                                   mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                st.download_button('Download PowerPoint (charts as images)', data=pptx_io,
                                   file_name='control_charts.pptx',
                                   mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')

                st.markdown("---")
                st.caption('Footer: OMAC Developer by SM Baqir 2025')

else:
    st.info('Upload an Excel workbook to get started.')
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

# --- Streamlit Page Setup ---
st.set_page_config(page_title="Control Chart Maker (QA OMAC Tools)", layout="wide")
st.markdown("<h1 style='text-align:left;'>Control chart maker (QA OMAC Tools)</h1>", unsafe_allow_html=True)
st.markdown("---")

# --- Sidebar ---
st.sidebar.header("Upload & Setup")
uploaded_file = st.sidebar.file_uploader("Upload an Excel workbook (.xlsx)", type=["xlsx"]) 

if uploaded_file is not None:
    try:
        xl = pd.ExcelFile(uploaded_file)
        sheets = xl.sheet_names
    except Exception as e:
        st.sidebar.error(f"Unable to read workbook: {e}")
        st.stop()

    sheet_name = st.sidebar.selectbox("Select sheet that contains data", options=sheets)
    if sheet_name:
        df = pd.read_excel(xl, sheet_name=sheet_name, header=0)
        st.write("**Preview of selected sheet**")
        st.dataframe(df.head())

        # --- First column is x-axis ---
        time_col = df.columns[0]
        st.sidebar.write(f"Time / Batch column detected: **{time_col}**")

        # --- Select parameter columns ---
        st.sidebar.subheader("Parameters for Control Chart")
        param_cols = st.sidebar.multiselect(
            "Select one or more parameter columns",
            options=[c for c in df.columns if c != time_col],
            default=[c for c in df.columns if c != time_col]
        )

        if len(param_cols) == 0:
            st.info("Select at least one parameter column from the sidebar to proceed.")
        else:
            if st.sidebar.button("Process & Generate Outputs"):
                working = df.copy()

                # --- Safe numeric conversion ---
                for col in param_cols:
                    working[col] = (
                        working[col]
                        .astype(str)
                        .str.replace(',', '', regex=False)
                        .str.strip()
                        .replace(['-', '–', '—', 'N/A', 'na', '', None], np.nan)
                    )
                    working[col] = pd.to_numeric(working[col], errors='coerce')

                # --- Prepare Excel output ---
                out_xlsx = BytesIO()
                writer = pd.ExcelWriter(out_xlsx, engine='openpyxl')

                # --- Prepare PPTX ---
                prs = Presentation()
                slide_layout = prs.slide_layouts[5]

                # --- Process each parameter ---
                for col in param_cols:
                    col_values = working[col].dropna()
                    if col_values.empty:
                        st.warning(f"No numeric data found for '{col}'. Skipping.")
                        continue

                    # Calculate statistics
                    CL = col_values.mean()
                    MR = col_values.diff().abs()
                    MRbar = MR[1:].mean() if len(MR) > 1 else 0
                    d2 = 1.128
                    UCL = CL + 3 * (MRbar / d2)
                    LCL = max(0, CL - 3 * (MRbar / d2))

                    # Append new columns
                    working[f'{col}_CL'] = CL
                    working[f'{col}_MR'] = MR
                    working[f'{col}_MRbar'] = MRbar
                    working[f'{col}_UCL'] = UCL
                    working[f'{col}_LCL'] = LCL

                    # --- Create chart ---
                    fig, ax = plt.subplots(figsize=(10,4))
                    ax.plot(working[time_col], col_values, marker='o', label=col)
                    ax.axhline(CL, linestyle='--', label='CL')
                    ax.axhline(UCL, color='r', linestyle='--', label='UCL')
                    ax.axhline(LCL, color='r', linestyle='--', label='LCL')
                    ax.set_title(f'Control Chart - {col}')
                    ax.set_xlabel(time_col)
                    ax.set_ylabel(col)
                    ax.legend()
                    plt.xticks(rotation=45)
                    plt.tight_layout()

                    # Save chart as image
                    img_bytes = BytesIO()
                    fig.savefig(img_bytes, format='png', dpi=150)
                    img_bytes.seek(0)
                    plt.close(fig)

                    # Add slide
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(img_bytes, Inches(0.5), Inches(0.7), width=Inches(9))

                # --- Save Excel ---
                working.to_excel(writer, sheet_name=sheet_name, index=False)
                writer.close()
                out_xlsx.seek(0)

                # --- Save PPTX ---
                pptx_io = BytesIO()
                prs.save(pptx_io)
                pptx_io.seek(0)

                # --- Streamlit download buttons ---
                st.success('Processing complete!')
                st.download_button('Download updated Excel', data=out_xlsx,
                                   file_name='control_chart_output.xlsx',
                                   mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                st.download_button('Download PowerPoint (charts as images)', data=pptx_io,
                                   file_name='control_charts.pptx',
                                   mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')

                st.markdown("---")
                st.caption('Footer: OMAC Developer by SM Baqir 2025')

else:
    st.info('Upload an Excel workbook to get started.')
